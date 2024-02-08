from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np

if __name__ == "__main__":
    print("Let's go!")

    # Creating presentation
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    picture_and_text_slide_layout = prs.slide_layouts[8]

    # Creating title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # Add Title
    title.text = "This is a great title"
    # Add Subtitle
    subtitle.text = 'With a punchline'

    # Creating a slide with bulletpoints
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Here comes the fun part"

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = 'This is a list'
    p = tf.add_paragraph()
    p.text = 'This is an item'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'This is a sub-item'
    p.level = 2

    # Adding a slide with a table
    slide = prs.slides.add_slide(picture_and_text_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'Adding a Table'

    df = pd.read_csv('example.csv')
    cols = len(df.columns.to_list())
    rows = len(df) + 1
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # Creating header
    for index, col in enumerate(df.columns.to_list()):
        table.cell(0, index).text = col

    # filling rows
    for row_index, row in df.iterrows():
        for col_index, col_name in enumerate(row):
            table.cell(row_index + 1, col_index).text = str(col_name)

    # Saving the Presentation
    prs.save(r"asd.pptx")
